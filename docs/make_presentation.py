"""Generate docs/RAG_Presentation.pptx — a teaching deck for the rag-simple project.

Usage:
    python docs/make_presentation.py

Re-run after editing the slide content below to regenerate the file.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT_FILE = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / "RAG_Presentation.pptx"

# --- Theme ---
DARK = RGBColor(0x1E, 0x25, 0x33)      # slide background
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)    # light blue
TEXT = RGBColor(0xEC, 0xEF, 0xF4)      # near-white
MUTED = RGBColor(0x9A, 0xA5, 0xB5)     # gray
CODE_BG = RGBColor(0x12, 0x17, 0x21)   # code block background
CODE_FG = RGBColor(0xA5, 0xD6, 0xA7)   # code text (soft green)
BODY_FONT = "Segoe UI"
CODE_FONT = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, lines, *, size=18, color=TEXT,
             font=BODY_FONT, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.15):
    """lines: list of (text, level) tuples or plain strings."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        indent = "    " * level
        bullet = "" if level == 0 else "– "
        run = p.add_run()
        run.text = f"{indent}{bullet}{text}" if level else text
        f = run.font
        f.name = font
        f.size = Pt(size - 2 * level)
        f.color.rgb = color if level == 0 else MUTED
        f.bold = bold and level == 0
    return box


def add_code(slide, left, top, width, height, code, *, size=14):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0x33, 0x3D, 0x4F)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    for i, line in enumerate(code.rstrip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        f = run.font
        f.name = CODE_FONT
        f.size = Pt(size)
        f.color.rgb = CODE_FG
    return box


def title_bar(slide, title, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9),
             [title], size=32, color=ACCENT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.5),
                 [subtitle], size=16, color=MUTED)
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.55), Inches(12.1), Emu(20000))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


def bullets_slide(title, subtitle, lines, code=None, code_size=13):
    slide = add_slide()
    title_bar(slide, title, subtitle)
    if code:
        add_text(slide, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.3), lines, size=17)
        add_code(slide, Inches(6.7), Inches(1.8), Inches(6.0), Inches(5.2), code, size=code_size)
    else:
        add_text(slide, Inches(0.6), Inches(1.9), Inches(12.1), Inches(5.2), lines, size=19)
    return slide


# ---------------------------------------------------------------- Slide 1: Title
slide = add_slide()
add_text(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1.4),
         ["Building RAG From Scratch"], size=54, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8),
         ["Retrieval-Augmented Generation with Python, NumPy, and Ollama"],
         size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.6), Inches(11.3), Inches(0.6),
         ["No frameworks · No API keys · ~250 lines of code · Runs on your laptop"],
         size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- Slide 2: Agenda
bullets_slide("Agenda", None, [
    "1.  Tech stack and prerequisites",
    "2.  The problem: what LLMs cannot do",
    "3.  The RAG idea in one sentence",
    "4.  Architecture: two phases, three files",
    "5.  Concept deep-dives: embeddings, cosine similarity, chunking",
    "6.  Code walkthrough: rag.py, ingest.py, ask.py — function by function",
    "7.  End-to-end trace of one question",
    "8.  Live demo commands",
    "9.  Exercises and where to go next",
])

# ---------------------------------------------------------------- Slide 3: Tech stack & prerequisites
bullets_slide("Tech Stack & Prerequisites",
              "Deliberately minimal — no LangChain, no cloud APIs, no database server", [
    "Tech stack:",
    ("Python 3.10+ — all code (~250 lines, 3 files)", 1),
    ("NumPy — embedding matrix, cosine math, argsort ranking", 1),
    ("Ollama — runs both models locally (localhost:11434)", 1),
    ("nomic-embed-text (~274 MB) — text → 768-dim vectors", 1),
    ("llama3.2 (3B, ~2 GB) — writes the grounded answer", 1),
    ("Flat files (.npy + .json) — our 'vector database'", 1),
    "",
    "Prerequisites:",
    ("8 GB RAM, ~3 GB disk, any OS — GPU NOT required", 1),
    ("Internet only for setup; afterwards fully offline", 1),
    ("Basic Python knowledge — no ML background needed", 1),
], code="""\
# 1. Install Ollama from https://ollama.com
ollama --version

# 2. Pull the two models
ollama pull nomic-embed-text   # embedder
ollama pull llama3.2           # chat model
ollama list                    # verify both

# 3. Python environment
python -m venv .venv
.venv\\Scripts\\activate        # Windows
source .venv/bin/activate      # macOS/Linux
pip install -r requirements.txt
# requirements.txt is 2 lines:
#   ollama>=0.4.0
#   numpy>=1.26.0

# 4. Verify
python ingest.py
python ask.py "Which planet is largest?"\
""", code_size=12)

# ---------------------------------------------------------------- Slide 3: Problem
bullets_slide("The Problem: LLMs Don't Know Your Data", None, [
    "An LLM is trained once, on a fixed snapshot of public text.",
    ("Knowledge cutoff — nothing after training day exists for the model", 1),
    ("Private data — it has never seen your notes, wikis, or PDFs", 1),
    ("Hallucination — when unsure, it invents confident-sounding answers", 1),
    "",
    "Fine-tuning to fix this is expensive, slow, and must be repeated for every update.",
    "",
    "We need a way to hand the model fresh, private knowledge at question time.",
])

# ---------------------------------------------------------------- Slide 4: RAG idea
bullets_slide("The RAG Idea in One Sentence", None, [
    "\u201cBefore asking the LLM, search your own documents for relevant passages",
    "and paste them into the prompt — then tell the model to answer ONLY from them.\u201d",
    "",
    "The model doesn't need to know the answer.",
    "It only needs to READ and SUMMARIZE the context we hand it.",
    "",
    "Retrieval  =  find the right passages   (search problem)",
    "Augmented  =  paste them into the prompt (string formatting!)",
    "Generation =  LLM writes the final answer (one API call)",
])

# ---------------------------------------------------------------- Slide 5: Architecture
slide = add_slide()
title_bar(slide, "Architecture: Two Phases, Three Files")
add_code(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(3.2), """\
PHASE 1 - INGEST (offline, run once)                       ingest.py
  data/*.txt --> chunk_text() --> embed() --> save_index()
                 500 chars,       768-dim        index/embeddings.npy
                 100 overlap      vectors        index/chunks.json

PHASE 2 - ASK (online, every question)                     ask.py
  question --> embed() --> retrieve() --> build_prompt() --> LLM
               768-dim     cosine sim,    context + rules    llama3.2
               vector      top-k chunks                      --> answer\
""", size=13)
add_text(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.6), [
    "Why two phases?  Embedding is slow (one model call per chunk).",
    "Do it once offline; each question then costs only 1 embed + 1 chat call.",
    "rag.py is the shared toolbox — both scripts import from it.",
], size=17)

# ---------------------------------------------------------------- Slide 6: Embeddings
bullets_slide("Concept 1 — Embeddings: Meaning as Numbers",
              "rag.embed()  ·  model: nomic-embed-text", [
    "An embedding model converts text into a fixed-length vector.",
    ("Here: every text becomes exactly 768 floating-point numbers", 1),
    "",
    "Key property: similar meaning ⇒ similar direction —",
    "even with zero shared words.",
    "",
    "The same function embeds document chunks AND questions,",
    "so both live in the same vector space and can be compared.",
], code="""\
def embed(text):
    response = ollama.embeddings(
        model=EMBED_MODEL,   # nomic-embed-text
        prompt=text)
    return np.array(response["embedding"],
                    dtype=np.float32)

"Jupiter is the largest planet"
  -> [ 0.12, -0.83, 0.44, ...]   768 dims
"The biggest world orbiting the Sun"
  -> [ 0.11, -0.79, 0.41, ...]   CLOSE!
"How to bake chocolate cake"
  -> [-0.55,  0.20, -0.91, ...]  far away\
""")

# ---------------------------------------------------------------- Slide 7: Cosine similarity
bullets_slide("Concept 2 — Cosine Similarity",
              "rag.cosine_similarity()", [
    "Measures the ANGLE between two vectors:",
    ("1.0 → same direction → same meaning", 1),
    ("0.0 → perpendicular → unrelated", 1),
    "",
    "Dividing by the norms means only direction matters,",
    "not text length.",
    "",
    "In this project a good match scores ≈ 0.55–0.80;",
    "irrelevant chunks fall below ≈ 0.4.",
], code="""\
def cosine_similarity(a, b):
    return np.dot(a, b) / (
        np.linalg.norm(a) * np.linalg.norm(b))

# Worked example in 2-D:
a = [1, 0];  b = [1, 1]
dot(a, b) = 1*1 + 0*1 = 1
|a| = 1,  |b| = sqrt(2) = 1.414
similarity = 1 / 1.414 = 0.707   (45 deg)\
""")

# ---------------------------------------------------------------- Slide 8: Chunking
bullets_slide("Concept 3 — Chunking With Overlap",
              "rag.chunk_text()  ·  CHUNK_SIZE=500, OVERLAP=100", [
    "Why not embed whole documents?",
    ("A 5-page doc's vector is a blurry average of all its topics", 1),
    ("Small chunks have sharp meanings that match questions", 1),
    ("Retrieved chunks must also FIT in the LLM prompt", 1),
    "",
    "Sliding window: advance by 500 − 100 = 400 chars,",
    "so consecutive chunks share 100 characters.",
    "",
    "Overlap ⇒ a sentence cut at a boundary still",
    "appears complete in at least one chunk.",
], code="""\
def chunk_text(text, chunk_size=500, overlap=100):
    chunks, start = [], 0
    while start < len(text):
        chunk = text[start:start + chunk_size]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap
    return chunks

Position: 0      400  500      800 900  1000
Chunk 1:  [======== 0-500 ====]
Chunk 2:         [====== 400-900 ====]
Chunk 3:                      [= 800-1000 =]
                 ^^overlap^^  ^^overlap^^\
""")

# ---------------------------------------------------------------- Slide 9: ingest.py
bullets_slide("Phase 1 — ingest.py, Step by Step",
              "python ingest.py", [
    "1.  LOAD — read every data/*.txt (sorted for reproducibility);",
    ("keep the filename with the text for citations later", 1),
    "2.  CHUNK — split each doc; record text, source, chunk_index",
    "3.  EMBED — one Ollama call per chunk (the slow part);",
    ("stack vectors into an N × 768 NumPy matrix", 1),
    "4.  SAVE — matrix → embeddings.npy, metadata → chunks.json",
    "",
    "Re-run whenever files in data/ change —",
    "the index does not update itself!",
], code="""\
documents = []
for fp in sorted(DATA_DIR.glob("*.txt")):
    documents.append({"source": fp.name,
                      "text": fp.read_text()})

all_chunks = []
for doc in documents:
    for i, c in enumerate(chunk_text(doc["text"])):
        all_chunks.append({"text": c,
                           "source": doc["source"],
                           "chunk_index": i})

embeddings = np.array(
    [embed(c["text"]) for c in all_chunks])
save_index(embeddings, all_chunks)

# Output: 28 chunks x 768 dimensions\
""")

# ---------------------------------------------------------------- Slide 10: The index
bullets_slide("The 'Vector Database' = Two Plain Files",
              "rag.save_index() / rag.load_index()", [
    "embeddings.npy — NumPy matrix, one 768-float row per chunk",
    "chunks.json — the chunk texts + source filename + position",
    "",
    "Row i of the matrix  ⟷  entry i of the JSON list.",
    "That positional pairing is the entire database schema!",
    "",
    "Real vector DBs (Chroma, Pinecone, pgvector) do exactly this,",
    "plus indexing tricks to search millions of rows fast.",
], code="""\
def save_index(embeddings, chunks):
    INDEX_DIR.mkdir(exist_ok=True)
    np.save(EMBEDDINGS_FILE, embeddings)
    CHUNKS_FILE.write_text(
        json.dumps(chunks, indent=2))

# chunks.json (one entry):
{
  "text": "Jupiter is the largest planet...",
  "source": "solar_system.txt",
  "chunk_index": 3
}\
""")

# ---------------------------------------------------------------- Slide 11: retrieve
bullets_slide("Phase 2a — retrieve(): The Heart of RAG",
              "rag.retrieve()", [
    "1.  Embed the question into the same 768-dim space",
    "2.  Score EVERY chunk with cosine similarity",
    ("brute force is instant for ~30 chunks", 1),
    "3.  argsort → ascending indices; [::-1] → descending;",
    ("[:top_k] → keep the k best (default 3)", 1),
    "4.  Return text + source + score for each winner",
    "",
    "At millions of chunks you'd swap the loop for a",
    "vector DB (FAISS, HNSW) — same concept, faster search.",
], code="""\
def retrieve(question, embeddings, chunks,
             top_k=3):
    q = embed(question)
    scores = [cosine_similarity(q, e)
              for e in embeddings]
    top = np.argsort(scores)[::-1][:top_k]
    return [{"text":   chunks[i]["text"],
             "source": chunks[i]["source"],
             "score":  float(scores[i])}
            for i in top]

# argsort([0.2, 0.9, 0.5]) -> [0, 2, 1]
# [::-1]                   -> [1, 2, 0]
# [:2]                     -> [1, 2]\
""")

# ---------------------------------------------------------------- Slide 12: prompt
bullets_slide("Phase 2b — build_prompt(): Grounding the LLM",
              "rag.build_prompt()", [
    "Three defenses against hallucination:",
    "",
    "1.  \u201cAnswer using ONLY the context below\u201d",
    ("blocks the model's general training knowledge", 1),
    "2.  An explicit escape hatch",
    ("the model is TOLD what to say when context is insufficient", 1),
    "3.  \u201cMention which source document(s)\u201d",
    ("every chunk is labeled [Source: file] — answers cite it", 1),
], code="""\
You are a helpful teaching assistant.
Answer the question using ONLY the
context below.
If the context does not contain enough
information, say "I don't have enough
information to answer that question."
Always mention which source document(s)
your answer comes from.

CONTEXT:
[Source: solar_system.txt]
Jupiter is the largest planet...
---
[Source: solar_system.txt]
Saturn is famous for its rings...

QUESTION: Which planet is the largest?
ANSWER:\
""")

# ---------------------------------------------------------------- Slide 13: generate
bullets_slide("Phase 2c — generate_answer(): The Smallest Function",
              "rag.generate_answer()", [
    "Build the grounded prompt, send one chat call, return the text.",
    "",
    "The \u201cG\u201d in RAG is the SIMPLEST part of the whole system.",
    "All the intelligence went into retrieval and prompt design.",
    "",
    "Central lesson of this project:",
    "RAG quality = retrieval quality + prompt quality.",
    "The LLM just reads what you found.",
], code="""\
def generate_answer(question, context_chunks):
    prompt = build_prompt(question,
                          context_chunks)
    response = ollama.chat(
        model=CHAT_MODEL,   # llama3.2
        messages=[{"role": "user",
                   "content": prompt}])
    return response["message"]["content"]\
""")

# ---------------------------------------------------------------- Slide 14: trace
slide = add_slide()
title_bar(slide, "End-to-End Trace of One Question",
          'python ask.py "Which planet is the largest?"')
add_code(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.3), """\
1. check_ollama()      server up, models installed
2. load_index()        28 x 768 matrix + 28 chunk records
3. embed(question)     "Which planet..." -> 768-dim vector
4. cosine sim x 28     Jupiter chunk 0.74 | Saturn 0.68 | photosynthesis 0.31
5. top-3               three solar_system.txt chunks win
6. build_prompt()      context + rules + question
7. ollama.chat()       llama3.2 reads context, writes answer
8. Output:             "Jupiter is the largest planet in our solar system...
                        (Source: solar_system.txt)"\
""", size=14)
add_text(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.9), [
    'Control case — "What is the capital of France?" → low scores, empty context → '
    '"I don\'t have enough information."  That is RAG SUCCEEDING: it refused to hallucinate.',
], size=15, color=ACCENT)

# ---------------------------------------------------------------- Slide 15: demo
bullets_slide("Live Demo Commands", None, [
    "Setup (once):",
], code="""\
ollama pull nomic-embed-text        # embedder, ~274 MB
ollama pull llama3.2                # chat model, ~2 GB
pip install -r requirements.txt

python ingest.py                    # Phase 1: build index

python ask.py "Which planet is the largest?"
python ask.py                       # interactive mode
python ask.py "Tell me about computers" --top-k 5
python ask.py "What is photosynthesis?" --quiet

# Watch the retrieval scores printed before each answer --
# they are your #1 debugging tool.\
""", code_size=14)

# ---------------------------------------------------------------- Slide 16: exercises
bullets_slide("Exercises", "From easy to challenging", [
    "1.  Add your own .txt to data/, re-ingest, ask about it",
    "2.  Set CHUNK_SIZE to 100, then 2000 — what breaks each way?",
    "3.  Compare --top-k 1 vs --top-k 5 on the same question",
    "4.  Find a score threshold that separates good from bad questions",
    "5.  Swap the chat model (ollama pull mistral; CHAT_MODEL = \"mistral\")",
    "6.  Challenge: rewrite chunk_text() to split on sentence boundaries",
    "7.  Challenge: delete the \u201cONLY the context\u201d rule and watch the",
    ("model hallucinate about France — grounding matters!", 1),
])

# ---------------------------------------------------------------- Slide 17: summary
bullets_slide("Summary — What You Now Know", None, [
    "RAG = Retrieve relevant chunks → Augment the prompt → Generate grounded answer",
    "",
    "Embeddings turn meaning into vectors; cosine similarity compares them",
    "Chunking with overlap gives sharp, prompt-sized search units",
    "An entire 'vector database' can be a .npy matrix + a .json list",
    "Retrieval is argsort over similarity scores — 3 lines of NumPy",
    "Prompt design (ONLY-context rule + escape hatch + citations) prevents hallucination",
    "",
    "Read the code in this order:  rag.py  →  ingest.py  →  ask.py",
    "Full walkthrough: docs/STUDENT_GUIDE.md",
])

prs.save(OUT_FILE)
print(f"Saved {OUT_FILE} ({len(prs.slides._sldIdLst)} slides)")
